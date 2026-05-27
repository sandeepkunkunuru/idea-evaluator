"""Absorb an investor pitch deck (.pptx) into the evaluator schema.

Pipeline:
  1. Extract per-slide text via python-pptx
  2. Deterministically pull numeric anchors (revenue, ask, TAM/SAM, headcount,
     IRR, growth pct, etc.) using money/percent regex over the full corpus
  3. If the Claude CLI is available, ask it to fill the remaining qualitative
     and structured fields. Otherwise emit what we have + defaults.

Output: a dict shaped like `schema.defaults()` — directly consumable by
`cli.py --input <generated>.json`.
"""
from __future__ import annotations

import html
import json
import re
import shutil
import subprocess
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from . import schema


# ---------- 1. Slide text extraction ----------------------------------------

def extract_slides(path: str) -> List[Dict]:
    """Return [{idx, title, lines:[...]}] for every slide."""
    from pptx import Presentation
    prs = Presentation(path)
    slides: List[Dict] = []
    for i, s in enumerate(prs.slides, 1):
        lines: List[str] = []
        title = ""
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs).strip()
                if not t:
                    continue
                # Some decks (e.g. PptxGenJS-generated) double-encode entities,
                # so run text arrives as "&amp;" / "&lt;". Decode to real chars.
                t = html.unescape(t)
                lines.append(t)
                if not title and sh == s.shapes.title if s.shapes.title else False:
                    title = t
        if not title and lines:
            title = lines[0]
        slides.append({"idx": i, "title": title, "lines": lines})
    return slides


# ---------- 2. Deterministic numeric extraction -----------------------------

# India crore/lakh + USD/EUR + bare millions/billions
_MONEY_RE = re.compile(
    r"""(?ix)
    (?P<sign>[₹$€£]|Rs\.?|INR|USD|EUR|GBP)?      # optional currency
    \s*
    (?P<num>\d{1,3}(?:[,\s]\d{2,3})*(?:\.\d+)?|\d+(?:\.\d+)?)
    \s*
    (?P<unit>cr(?:ore)?s?|lakhs?|lacs?|million|billion|mn|bn|k|m|b)?
    """
)

_PCT_RE = re.compile(r"(\d+(?:\.\d+)?)\s*%")

_USD_PER_INR = 1 / 83.0  # rough; only used when the deck speaks crores

_UNIT_MULT = {
    "cr": 1e7, "crore": 1e7, "crores": 1e7,
    "lakh": 1e5, "lakhs": 1e5, "lac": 1e5, "lacs": 1e5,
    "million": 1e6, "mn": 1e6, "m": 1e6,
    "billion": 1e9, "bn": 1e9, "b": 1e9,
    "k": 1e3,
}


def _to_usd(match: re.Match) -> Optional[float]:
    raw = match.group("num").replace(",", "").replace(" ", "")
    try:
        n = float(raw)
    except ValueError:
        return None
    unit = (match.group("unit") or "").lower()
    sign = (match.group("sign") or "").upper()
    # Reject bare numbers with no currency and no unit — these are almost always
    # counts (slide indexes, KPIs like "704/mo") not money.
    if not unit and not sign:
        return None
    mult = _UNIT_MULT.get(unit, 1.0)
    val = n * mult
    if sign in ("₹", "RS", "RS.", "INR"):
        val *= _USD_PER_INR
    elif sign in ("€", "EUR"):
        val *= 1.07
    elif sign in ("£", "GBP"):
        val *= 1.27
    return val


def _largest_money(blocks: List[str], *needles: str) -> Optional[float]:
    """Find the largest money figure in any block (slide) whose text mentions
    one of `needles`. Slides are passed in as blocks so adjacent stat-card text
    (e.g. "₹8 Cr+" + "Revenue Achieved") counts as co-located."""
    best = 0.0
    for blk in blocks:
        low = blk.lower()
        if needles and not any(n in low for n in needles):
            continue
        for m in _MONEY_RE.finditer(blk):
            v = _to_usd(m)
            if v and v > best:
                best = v
    return best or None


def _pct_near(blocks: List[str], *needles: str) -> Optional[float]:
    for blk in blocks:
        low = blk.lower()
        if needles and not any(n in low for n in needles):
            continue
        m = _PCT_RE.search(blk)
        if m:
            return float(m.group(1)) / 100.0
    return None


def deterministic_anchors(slides: List[Dict]) -> Dict:
    """Pull numeric anchors that we can extract reliably without an LLM."""
    blocks = ["\n".join(s["lines"]) for s in slides]
    out: Dict = {}

    ask = _largest_money(blocks, "ask", "investment", "raise", "funding", "equity")
    if ask:
        out["_ask_usd"] = ask

    rev = _largest_money(blocks, "revenue", "topline", "sales", "turnover")
    if rev:
        out["rev_y1"] = rev / 5  # if the deck shows 5-yr cumulative, rough divide

    tam = _largest_money(blocks, "tam", "addressable", "market size", "opportunity")
    if tam:
        out["tam_usd"] = tam
        out["sam_usd"] = tam * 0.1
        out["som_usd"] = tam * 0.01

    irr = _pct_near(blocks, "irr")
    if irr:
        out["_irr"] = irr
    growth = _pct_near(blocks, "cagr", "growth")
    if growth:
        out["market_growth_pct"] = growth
    margin = _pct_near(blocks, "ebitda margin", "gross margin", "pat margin")
    if margin:
        out["gross_margin"] = margin

    # Title -> idea_name
    if slides and slides[0]["lines"]:
        out["idea_name"] = slides[0]["lines"][0][:80]
        if len(slides[0]["lines"]) > 1:
            out["one_liner"] = slides[0]["lines"][1][:140]

    return out


# ---------- 3. LLM-assisted mapping -----------------------------------------

_FIELD_DOC = "\n".join(
    f"  - {f.key} ({f.kind}): {f.prompt}"
    for f in schema.ALL_FIELDS
    if not f.key.startswith("score_")
       and f.section not in ("equity",)
)

_LLM_PROMPT = """You will read a pitch deck and emit a JSON object that maps it
into a structured schema for valuation and rubric scoring. Output ONLY a JSON
object — no markdown, no commentary.

For monetary fields output USD numbers (convert ₹/Cr/Lakh if present:
1 Cr = 10,000,000 INR ≈ 120,500 USD; 1 Lakh = 100,000 INR ≈ 1,205 USD).
For percentages output 0..1 (e.g. 0.18 for 18%).
For 1-5 rubric scores, infer from how strongly the deck supports each claim.
Omit fields you genuinely cannot infer.

Schema fields you may set:
{fields}

Also include a top-level "_notes" string summarizing what the deck is about
in <= 3 sentences.

DECK (slide-by-slide):
{deck}
"""


def llm_map(slides: List[Dict], timeout: int = 180) -> Optional[Dict]:
    if not shutil.which("claude"):
        return None
    deck = "\n\n".join(
        f"--- Slide {s['idx']}: {s['title']} ---\n" + "\n".join(s["lines"])
        for s in slides
    )
    prompt = _LLM_PROMPT.format(fields=_FIELD_DOC, deck=deck)
    try:
        proc = subprocess.run(
            ["claude", "-p", prompt],
            capture_output=True, text=True, timeout=timeout,
        )
    except (subprocess.TimeoutExpired, OSError):
        return None
    if proc.returncode != 0:
        return None
    raw = proc.stdout.strip()
    # Defensive: strip code-fence wrappers if the model added them.
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)
    # Grab first balanced {...}
    start = raw.find("{")
    if start < 0:
        return None
    depth = 0
    for i, ch in enumerate(raw[start:], start):
        if ch == "{":
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0:
                blob = raw[start:i + 1]
                try:
                    return json.loads(blob)
                except json.JSONDecodeError:
                    return None
    return None


# ---------- 4. Top-level entry point ----------------------------------------

def absorb(path: str, use_llm: bool = True) -> Tuple[Dict, Dict]:
    """Return (inputs_dict, diagnostics). inputs_dict is ready for cli.py."""
    slides = extract_slides(path)
    anchors = deterministic_anchors(slides)

    merged = schema.defaults()
    # Apply anchors first (cheap, deterministic)
    for k, v in anchors.items():
        if k in schema.FIELDS_BY_KEY:
            merged[k] = schema.coerce(schema.FIELDS_BY_KEY[k], v)
        elif k.startswith("_"):
            merged[k] = v

    llm_out: Dict = {}
    if use_llm:
        llm_out = llm_map(slides) or {}
        for k, v in llm_out.items():
            if k in schema.FIELDS_BY_KEY:
                try:
                    merged[k] = schema.coerce(schema.FIELDS_BY_KEY[k], v)
                except Exception:
                    pass
            elif k.startswith("_"):
                merged[k] = v

    diag = {
        "slide_count": len(slides),
        "anchors_found": sorted(anchors.keys()),
        "llm_used": bool(llm_out),
        "llm_keys": sorted(llm_out.keys()) if llm_out else [],
    }
    return merged, diag
