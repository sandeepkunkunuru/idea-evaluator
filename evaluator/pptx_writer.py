"""Produce an investor pitch deck (.pptx) from the evaluator's outputs.

Slide structure modeled on a standard Seed/Series-A deck:
  1. Cover
  2. The Opportunity (problem + market stats)
  3. Solution / Why us
  4. Traction / current performance
  5. Business model & revenue snapshot
  6. Roadmap (org headcount growth as proxy)
  7. Market strategy / GTM
  8. Financial projections (DCF + comps)
  9. Valuation summary (all methods)
 10. Team & AI-native org
 11. Founder equity split
 12. Investment ask & use of funds
 13. Closing

Everything is generated from the same `inputs + rubric + valuations + equity`
dicts that drive the xlsx, so the deck always tells the same story.
"""
from __future__ import annotations

from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


# Palette — calm, investor-friendly
NAVY = RGBColor(0x0F, 0x2E, 0x5C)
TEAL = RGBColor(0x0E, 0x8E, 0x8E)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
GREY = RGBColor(0x4A, 0x4A, 0x4A)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _fmt_money(usd: float, prefer_inr_crore: bool = False) -> str:
    if not usd:
        return "—"
    if prefer_inr_crore:
        cr = usd * 83 / 1e7
        return f"₹{cr:,.1f} Cr"
    if usd >= 1e9:
        return f"${usd / 1e9:,.2f}B"
    if usd >= 1e6:
        return f"${usd / 1e6:,.2f}M"
    if usd >= 1e3:
        return f"${usd / 1e3:,.0f}K"
    return f"${usd:,.0f}"


def _add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _add_rect(slide, x, y, w, h, fill: RGBColor):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _add_text(slide, x, y, w, h, text, *, size=18, bold=False,
              color: RGBColor = GREY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def _slide_header(slide, title: str, subtitle: Optional[str] = None):
    _add_rect(slide, 0, 0, Inches(13.333), Inches(0.9), NAVY)
    _add_text(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.6),
              title, size=26, bold=True, color=WHITE)
    if subtitle:
        _add_text(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
                  subtitle, size=14, color=GREY)


def _stat_card(slide, x, y, w, h, big, label, color: RGBColor = TEAL):
    _add_rect(slide, x, y, w, h, LIGHT)
    _add_text(slide, x, y + Emu(60_000), w, Inches(0.9), big,
              size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    _add_text(slide, x, y + Inches(1.1), w, Inches(0.6), label,
              size=12, color=GREY, align=PP_ALIGN.CENTER)


def _bullets(slide, x, y, w, h, items: List[str], *, size=16):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = GREY
        p.space_after = Pt(6)


# ---------- Slide builders --------------------------------------------------

def _cover(prs, inputs):
    s = _add_blank(prs)
    _add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    _add_text(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.5),
              (inputs.get("idea_name") or "Untitled").upper(),
              size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.8),
              inputs.get("one_liner") or "",
              size=22, color=GOLD, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
              f"Investor Pitch — Stage: {inputs.get('stage','idea').title()} · "
              f"Geo: {inputs.get('geography','—')}",
              size=14, color=WHITE, align=PP_ALIGN.CENTER)


def _opportunity(prs, inputs):
    s = _add_blank(prs)
    _slide_header(s, "The Opportunity",
                  "Market size, growth, and the gap we close")
    tam = inputs.get("tam_usd", 0)
    sam = inputs.get("sam_usd", 0)
    som = inputs.get("som_usd", 0)
    growth = inputs.get("market_growth_pct", 0)
    _stat_card(s, Inches(0.5), Inches(1.6), Inches(3.0), Inches(2.0),
               _fmt_money(tam), "TAM")
    _stat_card(s, Inches(3.7), Inches(1.6), Inches(3.0), Inches(2.0),
               _fmt_money(sam), "SAM", color=GOLD)
    _stat_card(s, Inches(6.9), Inches(1.6), Inches(3.0), Inches(2.0),
               _fmt_money(som), "SOM (Yr-5)")
    _stat_card(s, Inches(10.1), Inches(1.6), Inches(2.7), Inches(2.0),
               f"{growth*100:,.0f}%", "Market CAGR", color=NAVY)
    _add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.5),
              "Problem", size=18, bold=True, color=NAVY)
    _add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
              inputs.get("problem") or "—", size=14, color=GREY)


def _solution(prs, inputs, rb):
    s = _add_blank(prs)
    _slide_header(s, "Our Solution",
                  "What we build, and why we win")
    _add_text(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(2.0),
              inputs.get("solution") or "—", size=14, color=GREY)
    # Strengths from rubric
    _add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
              "Top Strengths (rubric-scored)",
              size=18, bold=True, color=NAVY)
    items = [f"{k.replace('score_','').replace('_',' ').title()}: {v:.1f}/5  ({cat})"
             for k, v, cat in rb.get("strengths", [])[:5]]
    _bullets(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8), items or ["—"])


def _traction(prs, inputs):
    s = _add_blank(prs)
    _slide_header(s, "Traction & Current Performance")
    _stat_card(s, Inches(0.5), Inches(1.6), Inches(3.0), Inches(2.0),
               inputs.get("stage", "—").title(), "Stage")
    _stat_card(s, Inches(3.7), Inches(1.6), Inches(3.0), Inches(2.0),
               _fmt_money(inputs.get("rev_y1", 0)), "Projected Y1 Revenue",
               color=GOLD)
    _stat_card(s, Inches(6.9), Inches(1.6), Inches(3.0), Inches(2.0),
               f"{inputs.get('gross_margin',0)*100:,.0f}%", "Gross Margin")
    _stat_card(s, Inches(10.1), Inches(1.6), Inches(2.7), Inches(2.0),
               str(inputs.get("competitors_count", 0)),
               "Meaningful Competitors", color=NAVY)
    _add_text(s, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.5),
              "Target Customer (ICP)", size=18, bold=True, color=NAVY)
    _add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.0),
              inputs.get("target_customer") or "—", size=14, color=GREY)


def _projections(prs, inputs, vals):
    s = _add_blank(prs)
    _slide_header(s, "5-Year Financial Projections")
    rev = inputs.get("rev_y1", 0)
    g = [inputs.get(f"rev_growth_y{i}", 0) for i in (2, 3, 4, 5)]
    revs = [rev]
    for gr in g:
        revs.append(revs[-1] * (1 + gr))
    for i, r in enumerate(revs):
        x = Inches(0.5 + i * 2.5)
        _stat_card(s, x, Inches(1.6), Inches(2.3), Inches(2.0),
                   _fmt_money(r), f"Year {i+1}",
                   color=TEAL if i < 4 else GOLD)
    dcf = vals.get("dcf", {}).get("equity", 0)
    _add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
              f"DCF equity value (WACC {inputs.get('discount_rate',0)*100:.0f}%, "
              f"terminal g {inputs.get('terminal_growth',0)*100:.0f}%):  "
              f"{_fmt_money(dcf)}",
              size=18, bold=True, color=NAVY)


def _valuation(prs, vals):
    s = _add_blank(prs)
    _slide_header(s, "Valuation Across Seven Methods")
    methods = [
        ("DCF (equity)", vals.get("dcf", {}).get("equity", 0)),
        ("Comparables (mid)", vals.get("comparables", {}).get("midpoint", 0)),
        ("Berkus", vals.get("berkus", 0)),
        ("Scorecard (Payne)", vals.get("scorecard", 0)),
        ("VC method", vals.get("vc_method", 0)),
        ("First Chicago (wgt)", vals.get("first_chicago", {}).get("weighted", 0)),
        ("Patent (mid)", vals.get("patent", {}).get("midpoint", 0)),
    ]
    for i, (name, v) in enumerate(methods):
        col = i % 4
        row = i // 4
        x = Inches(0.5 + col * 3.2)
        y = Inches(1.5 + row * 2.3)
        _stat_card(s, x, y, Inches(3.0), Inches(2.0),
                   _fmt_money(v), name,
                   color=TEAL if i % 2 == 0 else GOLD)


def _team_org(prs, org_plan):
    s = _add_blank(prs)
    _slide_header(s, "AI-Native Org Structure",
                  "Future-proof titles, Y1 → Y3 ramp")
    rows = (org_plan if isinstance(org_plan, list) else org_plan.get("rows", []))[:14]
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.2),
                              Inches(12.3), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    hdr = tf.paragraphs[0]
    r = hdr.add_run()
    r.text = f"{'Role':<42} {'Y1':>4} {'Y3':>4} {'Automation':>11}"
    r.font.name = "Menlo"; r.font.size = Pt(12); r.font.bold = True
    r.font.color.rgb = NAVY
    for row in rows:
        p = tf.add_paragraph()
        rr = p.add_run()
        rr.text = (f"{row.get('role', row.get('title',''))[:42]:<42} "
                   f"{row.get('headcount_y1', row.get('y1',0)):>4} "
                   f"{row.get('headcount_y3', row.get('y3',0)):>4} "
                   f"{row.get('automation', row.get('automation_pct',0))*100:>10.0f}%")
        rr.font.name = "Menlo"; rr.font.size = Pt(11); rr.font.color.rgb = GREY


def _equity(prs, founders, esop_pct):
    s = _add_blank(prs)
    _slide_header(s, "Founder Equity Split",
                  "Slicing-Pie contribution units, net of ESOP")
    _add_text(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
              f"ESOP pool: {esop_pct*100:.0f}%",
              size=16, bold=True, color=GOLD)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.9),
                              Inches(12.3), Inches(5.0))
    tf = tb.text_frame; tf.word_wrap = True
    hdr = tf.paragraphs[0]
    r = hdr.add_run()
    r.text = f"{'Founder':<24} {'Units':>14} {'Equity %':>12}"
    r.font.name = "Menlo"; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = NAVY
    for f in founders:
        if f.get("units", 0) <= 0:
            continue
        p = tf.add_paragraph()
        rr = p.add_run()
        rr.text = (f"{str(f.get('name',''))[:24]:<24} "
                   f"{f.get('units',0):>14,.0f} "
                   f"{f.get('final_equity', f.get('pct_post_esop',0))*100:>11.1f}%")
        rr.font.name = "Menlo"; rr.font.size = Pt(12); rr.font.color.rgb = GREY


def _ask(prs, inputs, vals):
    s = _add_blank(prs)
    _slide_header(s, "Investment Ask")
    dcf = vals.get("dcf", {}).get("equity", 0) or 1
    ask_usd = inputs.get("_ask_usd") or dcf * 0.2  # default: 20% of DCF
    _stat_card(s, Inches(4.0), Inches(1.5), Inches(5.3), Inches(2.2),
               _fmt_money(ask_usd), "Equity Investment",
               color=GOLD)
    use = [
        ("Product & R&D", 0.45),
        ("Go-to-Market", 0.30),
        ("Team & Ops", 0.15),
        ("Working Capital", 0.10),
    ]
    for i, (k, p) in enumerate(use):
        x = Inches(0.5 + i * 3.2)
        _stat_card(s, x, Inches(4.0), Inches(3.0), Inches(2.0),
                   f"{int(p*100)}%", k,
                   color=TEAL if i % 2 == 0 else NAVY)


def _closing(prs, inputs):
    s = _add_blank(prs)
    _add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    _add_text(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.5),
              "Let's build it.",
              size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.8),
              inputs.get("idea_name", ""),
              size=20, color=WHITE, align=PP_ALIGN.CENTER)


# ---------- Public entry point ----------------------------------------------

def write_deck(inputs: Dict, results: Dict, path: str) -> None:
    """Generate a 10-slide pitch deck.

    `results` is the same dict the xlsx writer receives:
        {"rubric": ..., "valuations": ..., "org": ..., "equity": ...}
    Missing pieces are tolerated — slides degrade to "—".
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    rb = results.get("rubric", {})
    vals = results.get("valuations", {})
    org_plan = results.get("org", {"rows": []})
    eq = results.get("equity", {"founders": [], "esop_pct": inputs.get("esop_pool_pct", 0.15)})

    _cover(prs, inputs)
    _opportunity(prs, inputs)
    _solution(prs, inputs, rb)
    _traction(prs, inputs)
    _projections(prs, inputs, vals)
    _valuation(prs, vals)
    _team_org(prs, org_plan)
    _equity(prs, eq.get("founders", []), eq.get("esop_pct", 0.15))
    _ask(prs, inputs, vals)
    _closing(prs, inputs)

    prs.save(path)
